"""Generate PowerPoint presentation from the research data.

Usage: python scripts/generate_pptx.py
Requires: backend API running at localhost:8000
Output: outputs/presentation.pptx
"""

import json
import os
import sys
import requests
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

API = "http://localhost:8000"
OUTPUT = Path(__file__).parent.parent / "outputs" / "presentation.pptx"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
RED = RGBColor(0xE6, 0x19, 0x4B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE = RGBColor(0xF5, 0x82, 0x31)

YEAR_GROUPS = ['2015', '2016', '2017', '2018', '2019', '2020', '2021', '2022', '2023', '2024', '2025']

GROUP_COLORS_HEX = {
    'mrbeast': 'E6194B', '2015': '3CB44B', '2016': '4363D8', '2017': 'F58231',
    '2018': '911EB4', '2019': '42D4F4', '2020': 'F032E6', '2021': 'BFEF45',
    '2022': 'FABED4', '2023': '469990', '2024': 'DCBEFF', '2025': '000075',
}


def fetch(endpoint):
    r = requests.get(f"{API}{endpoint}", timeout=30)
    r.raise_for_status()
    return r.json()


# ─── Helpers ──────────────────────────────────────────────────

def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=Inches(0.4), font_size=32, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_subtitle(slide, text, top=Inches(1.2), font_size=16, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_bullets(slide, items, left=Inches(1), top=Inches(2), width=Inches(8), font_size=16):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.45 + 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.space_after = Pt(8)


def add_stat_card(slide, label, value, left, top, width=Inches(2), sub=None):
    shape = slide.shapes.add_shape(
        1, left, top, width, Inches(0.9)  # MSO_SHAPE.ROUNDED_RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLACK

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

    if sub:
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(8)
        p3.font.color.rgb = GRAY
        p3.alignment = PP_ALIGN.CENTER


def add_code_block(slide, code_text, left=Inches(0.5), top=Inches(3), width=Inches(4.5), height=Inches(2)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
        p.space_after = Pt(2)


def add_bar_chart(slide, categories, values, title="", left=Inches(0.5), top=Inches(2.2),
                  width=Inches(9), height=Inches(4), series_name="Value"):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    return chart


def add_line_chart(slide, categories, series_dict, left=Inches(0.5), top=Inches(2.2),
                   width=Inches(9), height=Inches(4)):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        chart_data.add_series(name, values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    if len(series_dict) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    return chart


def add_table(slide, headers, rows, left=Inches(1), top=Inches(2), col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = sum(col_widths) if col_widths else Inches(n_cols * 2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(n_rows * 0.4))
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    return table


# ─── Slide Builders ───────────────────────────────────────────

def slide_0_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_title(slide, "Clicking Toward Conformity", top=Inches(1.5), font_size=40)
    add_subtitle(slide, "Quantifying the Convergence of YouTube Thumbnail Design\nToward MrBeast's Visual Formula (2015–2025)", top=Inches(2.5), font_size=18)
    add_subtitle(slide, "Lucas Trinh  &  Zachary Chen", top=Inches(3.8), font_size=16, color=GRAY)
    add_subtitle(slide, "2026 HTCC Conference", top=Inches(4.4), font_size=14, color=GRAY)


def slide_1_questions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Research Questions")
    add_bullets(slide, [
        "Has YouTube thumbnail design converged toward MrBeast's visual style over 2015–2025?",
        "Which visual features changed most over time?",
        "Do modern thumbnails cluster closer to MrBeast's than historical ones?",
        "Is the shift statistically significant — or just noise?",
    ])


def slide_2_why_mrbeast(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Why MrBeast?")
    add_subtitle(slide, "Most subscribed individual creator (~300M+). Pioneered a distinct, widely-imitated thumbnail formula.")
    cards = [
        ("Have a Face", "81%", "vs 60% in 2015"),
        ("Are Smiling", "75%", "vs 47% in 2015"),
        ("Bright Thumbnails", "69%", "vs 41% in 2015"),
        ("Large Body in Frame", "53%", "vs 30% in 2015"),
    ]
    for i, (label, val, sub) in enumerate(cards):
        add_stat_card(slide, label, val, Inches(0.5 + i * 2.4), Inches(2.5), sub=sub)
    add_subtitle(slide, '"Big faces, big expressions, bright colors."', top=Inches(4.2), font_size=16, color=GRAY)


def slide_3_data_collection(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Data Collection")
    add_subtitle(slide, "YouTube Data API v3 — top 15 most-viewed videos per channel per year, 2015–2025")

    add_code_block(slide, """# YouTube API search by channel + year
response = youtube.search().list(
  channelId=channel_id,
  publishedAfter="2020-01-01",
  publishedBefore="2020-12-31",
  order="viewCount",
  maxResults=15)""", left=Inches(0.5), top=Inches(2), width=Inches(4.5), height=Inches(2.2))

    add_bullets(slide, [
        "15 thumbnails per channel per year (top by views, Shorts excluded)",
        "22 panel channels tracked across all 11 years",
        "309 MrBeast reference thumbnails (5 eras)",
        "Downloaded at max resolution (1280×720)",
    ], left=Inches(0.5), top=Inches(4.4), width=Inches(4.5), font_size=12)

    # Channel list on right side
    txBox = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    sections = [
        ("ENTERTAINMENT (14)", "Dude Perfect, PewDiePie, Sidemen, FaZe Rug, Danny Duncan, Logan Paul, KSI, IShowSpeed, Ryan Trahan, Airrack, JiDion, Unspeakable, David Dobrik, Matt Stonie"),
        ("GAMING (3)", "Markiplier, VanossGaming, Dream"),
        ("ART & OTHER (3)", "ZHC, Good Mythical Morning, Smosh"),
        ("CONTROLS (2)", "Kurzgesagt, CinemaSins"),
    ]

    p = tf.paragraphs[0]
    p.text = "Panel Channels (22)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK

    for cat, channels in sections:
        p = tf.add_paragraph()
        p.text = cat
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        p.text = channels
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)


def slide_4_dataset_overview(prs, overview):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Dataset Overview")
    cards = [
        ("Total Thumbnails", f"{overview['total_thumbnails']:,}"),
        ("Visual Features", "14"),
        ("Title Features", "9"),
        ("MrBeast Reference", "309"),
    ]
    for i, (label, val) in enumerate(cards):
        add_stat_card(slide, label, val, Inches(0.3 + i * 2.4), Inches(2))
    cards2 = [
        ("Panel Channels", "22", "Entertainment focus"),
        ("Time Span", "11 years", "2015–2025"),
        ("Features Extracted", f"{overview['features_extracted']:,}", None),
    ]
    for i, (label, val, sub) in enumerate(cards2):
        add_stat_card(slide, label, val, Inches(1 + i * 2.8), Inches(3.5), sub=sub)


def slide_feature(prs, title, subtitle, explanation, code, outputs, tool):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, f"Feature: {title}")
    add_subtitle(slide, subtitle)

    # Explanation
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = explanation
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)

    # Code
    add_code_block(slide, code, left=Inches(0.5), top=Inches(3.3), width=Inches(4.5), height=Inches(1.8))

    # Outputs
    txBox2 = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.5), Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Outputs"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK

    for out in outputs:
        p = tf2.add_paragraph()
        p.text = f"  {out}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
        p.space_after = Pt(4)

    p = tf2.add_paragraph()
    p.text = tool
    p.font.size = Pt(10)
    p.font.color.rgb = BLUE
    p.space_before = Pt(12)


def slide_11_scoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Scoring Systems")

    systems = [
        ("Weighted Likeness (Primary)", "Data-derived weights per feature based on discriminative power. smile (0.442), brightness (0.441), brow raise (0.270), body coverage (0.182)."),
        ("Continuous Similarity (0–100%)", "Z-score distance from MrBeast centroid across 10 discriminative features with exponential decay mapping."),
        ("Binary Likeness (0–8 pts, Reference)", "+1 for each threshold met: brightness ≥ 0.60, face count ≥ 1, text ≤ 0.005, smile ≥ 0.40, mouth open ≥ 0.15, body coverage ≥ 0.30, brow raise ≥ 0.30, face area ≥ 0.06"),
    ]
    for i, (name, desc) in enumerate(systems):
        y = Inches(1.8 + i * 1.5)
        txBox = slide.shapes.add_textbox(Inches(1), y, Inches(8), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY


def slide_12_formula(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "The MrBeast Formula vs. 2015 Baseline")

    headers = ["Feature", "MrBeast", "2015", "Change"]
    rows = [
        ["Face Count", "1.37", "0.78", "+75.6%"],
        ["Body Coverage", "0.413", "0.237", "+74.3%"],
        ["Smile Score", "0.443", "0.266", "+66.9%"],
        ["Mouth Open", "0.175", "0.105", "+66.7%"],
        ["Brow Raise", "0.332", "0.209", "+58.7%"],
        ["Largest Face Size", "0.087", "0.058", "+50.5%"],
        ["Brightness", "0.658", "0.570", "+15.5%"],
    ]
    add_table(slide, headers, rows, left=Inches(1.5), top=Inches(1.8),
              col_widths=[Inches(2), Inches(1.5), Inches(1.5), Inches(1.5)])


def slide_chart_bar(prs, title, subtitle, categories, values, series_name="Score",
                    stat_cards=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
    add_bar_chart(slide, categories, values, series_name=series_name)
    if stat_cards:
        for i, (label, val, sub) in enumerate(stat_cards):
            add_stat_card(slide, label, val, Inches(1 + i * 2.8), Inches(6.2), sub=sub)
    return slide


def slide_chart_line(prs, title, subtitle, categories, series_dict, stat_cards=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
    add_line_chart(slide, categories, series_dict)
    if stat_cards:
        for i, (label, val, sub) in enumerate(stat_cards):
            add_stat_card(slide, label, val, Inches(0.5 + i * 2.4), Inches(6.2), sub=sub)
    return slide


def slide_16_convergence(prs, face_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature-Level Convergence")
    add_bullets(slide, [
        "Face count: nearly closed the gap entirely",
        "Smile score: +18.9 percentage points",
        "Brow raise: 59–87% gap closure",
        "Body coverage: significant increase",
        "Mouth open score: +66.7%, matching MrBeast's expressive style",
    ], left=Inches(0.5), top=Inches(1.8), width=Inches(4.5))

    categories = YEAR_GROUPS
    values = [face_data['groups'].get(g, {}).get('mean', 0) for g in YEAR_GROUPS]
    add_bar_chart(slide, categories, values, title="Mean Face Count by Year",
                  left=Inches(5), top=Inches(1.8), width=Inches(4.8), height=Inches(3.5),
                  series_name="Face Count")


def slide_17_gap_closure(prs, gap_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Gap Closure: How Far Has the Industry Moved?")
    add_subtitle(slide, "Percentage of the 2015-to-MrBeast gap closed by 2025 (panel only)")

    features = [d['feature'] for d in gap_data]
    values = [d['closure'] for d in gap_data]

    chart_data = CategoryChartData()
    chart_data.categories = features
    chart_data.add_series("Gap Closed %", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2.2), Inches(9), Inches(4), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False


def slide_case_study(prs, channel_name, subtitle, years, scores, stat_cards, color=RED):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, f"Case Study: {channel_name}")
    add_subtitle(slide, subtitle)
    add_line_chart(slide, years, {channel_name: scores},
                   left=Inches(0.5), top=Inches(2), width=Inches(9), height=Inches(3.5))
    for i, (label, val, sub) in enumerate(stat_cards):
        add_stat_card(slide, label, val, Inches(0.3 + i * 2.4), Inches(5.8), sub=sub)


def slide_27_diffusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Diffusion of Innovation Model")

    phases = [
        ("Innovators", "2015–2017", "MrBeast pioneering the bright, face-forward, expressive formula", RED),
        ("Early Adopters", "2018–2019", "David Dobrik, ZHC, Ryan Trahan begin experimenting", ORANGE),
        ("Early Majority", "2020–2022", "Broad panel adoption; inflection point in convergence data", GREEN),
        ("Late Majority", "2023–2025", "Near-ubiquitous among entertainment channels", BLUE),
        ("Laggards / Resisters", "Ongoing", "PewDiePie (stable), Markiplier (diverging), genre-locked channels", PURPLE),
    ]
    for i, (phase, years, desc, color) in enumerate(phases):
        y = Inches(1.8 + i * 0.95)
        # Color bar
        bar = slide.shapes.add_shape(1, Inches(0.8), y, Inches(0.15), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # Text
        txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(8), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{phase}  ({years})"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY


def slide_28_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Limitations & Conclusions")

    add_bullets(slide, [
        "No engagement data (views/CTR not linked)",
        "Low R² (0.5–1.7%) — individual variation dominates",
        "Cannot prove causation — algorithm, tools, audience all plausible",
        "Panel selection bias toward entertainment",
    ], left=Inches(0.5), top=Inches(1.8), width=Inches(4.5), font_size=13)

    txBox = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(0.1), Inches(0.1))  # spacer

    add_bullets(slide, [
        "Convergence is statistically significant (p < 10⁻¹⁰, slope +0.105/yr)",
        "73% of panel channels converge (avg slope +0.127/yr)",
        "Multi-dimensional shift: face, smile, brow, body coverage",
        "Smile + brightness = 53% of discriminative weight",
    ], left=Inches(5.2), top=Inches(1.8), width=Inches(4.5), font_size=13)

    add_subtitle(slide, "Thank You", top=Inches(5.5), font_size=22, color=BLACK)
    add_subtitle(slide, "Lucas Trinh & Zachary Chen · 2026 HTCC Conference", top=Inches(6.2), font_size=14, color=GRAY)


# ─── Main ─────────────────────────────────────────────────────

def main():
    print("Fetching data from API...")
    try:
        overview = fetch("/stats/overview")
    except Exception as e:
        print(f"ERROR: Cannot connect to API at {API}. Is the backend running?\n{e}")
        sys.exit(1)

    weighted = fetch("/stats/weighted-likeness?panel_only=true")
    likeness = fetch("/stats/mrbeast-likeness?panel_only=true")
    similarity = fetch("/stats/mrbeast-similarity?panel_only=true")
    face_compare = fetch("/stats/compare?feature=face.face_count")
    convergence = fetch("/stats/convergence-tests?panel_only=true")
    evolution = fetch("/stats/channel-evolution?min_years=3&panel_only=true")
    weighted_weights = fetch("/stats/weighted-likeness?panel_only=true")
    title_data = fetch("/stats/title-likeness?panel_only=true")

    # Gap closure data
    gap_features = [
        ("face.face_count", "Face Count"),
        ("face.largest_face_area_ratio", "Face Size"),
        ("pose.body_coverage", "Body Coverage"),
        ("face.emotion_proxies.mouth_open_score", "Mouth Open"),
        ("face.emotion_proxies.smile_score", "Smile Score"),
        ("face.emotion_proxies.brow_raise_score", "Brow Raise"),
        ("color.avg_brightness", "Brightness"),
    ]
    gap_data = []
    for path, label in gap_features:
        d = fetch(f"/stats/compare?feature={path}")
        mb = d['groups'].get('mrbeast', {}).get('mean', 0)
        y15 = d['groups'].get('2015', {}).get('mean', 0)
        y25 = d['groups'].get('2025', {}).get('mean', 0)
        gap = mb - y15
        closure = round(((y25 - y15) / gap) * 100) if abs(gap) > 0.001 else 0
        gap_data.append({"feature": label, "closure": closure})
    gap_data.sort(key=lambda x: -x['closure'])

    print("Building presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 0: Title
    slide_0_title(prs)

    # Slide 1: Research Questions
    slide_1_questions(prs)

    # Slide 2: Why MrBeast?
    slide_2_why_mrbeast(prs)

    # Slide 3: Data Collection
    slide_3_data_collection(prs)

    # Slide 4: Dataset Overview
    slide_4_dataset_overview(prs, overview)

    # Slides 5-10: Feature Extraction
    features = [
        ("Color Analysis",
         "Converts each thumbnail to HSV color space to measure brightness, saturation, and color temperature.",
         "Each pixel's hue is classified as warm (reds, oranges: 0–30 and 150–179) or cool (greens, blues: 30–150). The ratio gives a single warm/cool score.",
         "# Warm/cool color scoring\nwarm = (hue <= 30) | (hue >= 150)\ncool = (hue > 30) & (hue < 150)\nscore = (warm - cool) / total",
         ["avg_saturation — mean color intensity", "avg_brightness — mean luminance", "warm_cool_score — -1 (cool) to +1 (warm)", "dominant_palette — top 5 hex colors", "hue_hist — 36-bin hue distribution"],
         "OpenCV, PIL, NumPy"),
        ("Face & Emotion Detection",
         "Uses MediaPipe FaceMesh (468 landmarks per face) to detect faces and approximate emotional expressions.",
         "Smile is measured by comparing mouth corner height to lip center height, normalized by mouth width. Higher values indicate upturned corners.",
         "# Smile from landmark geometry\nlip_y = (up_lip.y + lo_lip.y) / 2\ncorner_y = (left.y + right.y) / 2\nsmile = (lip_y - corner_y) / width\nsmile = clamp(smile + 0.5, 0, 1)",
         ["face_count — number of faces detected", "largest_face_area_ratio — biggest face as % of image", "smile_score — 0–1 smile intensity", "mouth_open_score — 0–1 mouth openness", "brow_raise_score — 0–1 eyebrow lift"],
         "MediaPipe FaceMesh"),
        ("Pose Detection",
         "Uses MediaPipe Pose (33 body landmarks) to measure how much of the frame a person occupies.",
         "Body coverage is the bounding box area of all visible pose landmarks divided by total image area. Higher values mean the subject dominates the thumbnail.",
         "# Body coverage from landmarks\nxs = [p.x for p in points]\nys = [p.y for p in points]\nw = max(xs) - min(xs)\nh = max(ys) - min(ys)\ncoverage = (w * h) / img_area",
         ["people_count — number of people detected", "body_coverage — 0–1 frame occupancy", "hand_visible_count — visible hands (0–2)", "pose_orientation — frontal, side, or back"],
         "MediaPipe Pose"),
        ("Text Detection",
         "Runs Tesseract OCR to detect and measure text overlays on thumbnails.",
         "OCR scans each thumbnail, filters results by confidence (>30%), then measures total text bounding box area as a fraction of the image.",
         "# OCR text detection\ndata = pytesseract.image_to_data(img)\nfor i in range(n):\n  if data[\"conf\"][i] > 30:\n    boxes.append(box)",
         ["has_text — boolean, any text detected", "text_area_ratio — text area / image area", "text_box_count — number of text regions", "detected_text — recognized strings"],
         "PyTesseract OCR"),
        ("Depth Estimation",
         "Uses MiDaS (monocular depth) to estimate foreground/background separation without stereo cameras.",
         "A pre-trained MiDaS model produces a per-pixel depth map from a single image. We then analyze the depth distribution.",
         "# MiDaS depth inference\nx = transform(img).to(device)\nwith torch.no_grad():\n  depth = model(x)\n  depth = F.interpolate(\n    depth, size=img.shape[:2],\n    mode=\"bicubic\")",
         ["depth_contrast — std of depth values", "foreground_ratio — % of pixels in foreground", "depth_range — max - min depth", "subject_depth_center — (x, y) of subject"],
         "MiDaS / PyTorch"),
        ("Title Analysis",
         "Parses video titles to detect MrBeast-style linguistic patterns using keyword matching and regex.",
         "Titles are tokenized, lowercased, then checked against curated word lists for money references, superlatives, and challenge framing.",
         "# Title pattern matching\nw = set(title.lower().split())\nmoney = re.search(r\"\\$[\\d,]+\", title)\nsuperlative = SUPERLATIVES & w\nchallenge = CHALLENGES & w",
         ["has_money_reference — dollar amounts detected", "has_superlative — extreme adjectives", "has_challenge_framing — competition language", "uppercase_ratio — ALL CAPS intensity", "word_count — title length"],
         "NLP / Regex"),
    ]
    for args in features:
        slide_feature(prs, *args)

    # Slide 11: Scoring Systems
    slide_11_scoring(prs)

    # Slide 12: MrBeast Formula
    slide_12_formula(prs)

    # Slide 13: Weighted Likeness Over Time
    w_values = [weighted['groups'].get(g, {}).get('normalized_mean', 0) * 100 for g in YEAR_GROUPS]
    slide_chart_bar(prs, "Weighted Likeness Over Time (Panel Only)",
                    "Normalized weighted score by year — captures gradual convergence",
                    YEAR_GROUPS, w_values, series_name="Weighted %",
                    stat_cards=[("Max Possible", f"{weighted['max_possible_score']:.2f}", "weighted points"),
                                ("MrBeast Mean", f"{weighted['groups'].get('mrbeast', {}).get('normalized_mean', 0) * 100:.1f}%", None)])

    # Slide 14: Binary Likeness
    b_values = [likeness['groups'].get(g, {}).get('mean_score', 0) for g in YEAR_GROUPS]
    slide_chart_bar(prs, "Binary Likeness Score (Reference)",
                    "Simple 0–8 threshold scoring",
                    YEAR_GROUPS, b_values, series_name="Mean Likeness",
                    stat_cards=[("2015 Baseline", "3.59", None), ("2024 Peak", "4.58", "+23% increase")])

    # Slide 15: Continuous Similarity
    s_values = [similarity['groups'].get(g, {}).get('mean_similarity', 0) for g in YEAR_GROUPS]
    slide_chart_line(prs, "Continuous Similarity Trend (Panel Only)",
                     "Mean z-score similarity to MrBeast centroid (0–100%)",
                     YEAR_GROUPS, {"Similarity %": s_values},
                     stat_cards=[("2015", "63.7%", None), ("2025", "67.9%", None)])

    # Slide 16: Feature-Level Convergence
    slide_16_convergence(prs, face_compare)

    # Slide 17: Gap Closure
    slide_17_gap_closure(prs, gap_data)

    # Slide 18: Clustering
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Clustering Analysis")
    add_subtitle(slide, "K-means clustering with PCA 2D projection (12 features). PCA explains 43.2% of variance.")
    cards = [
        ("Cluster 0 (Modern)", "59%", "of MrBeast here"),
        ("Cluster 1 (Classic)", "52%", "of 2015 here"),
        ("Centroid Distance", "1.05", "PCA Euclidean"),
        ("PCA Variance", "43.2%", "up from 33.4%"),
    ]
    for i, (label, val, sub) in enumerate(cards):
        add_stat_card(slide, label, val, Inches(0.3 + i * 2.4), Inches(3), sub=sub)

    # Slide 19: Channel Evolution
    trends = sorted(evolution['trends'], key=lambda x: -x['slope'])[:10]
    slide_chart_bar(prs, "Channel-Level Evolution",
                    "Per-channel likeness trends (panel channels, ≥3 years)",
                    [t['channel'] for t in trends],
                    [t['slope'] for t in trends],
                    series_name="Slope (pts/yr)",
                    stat_cards=[("Converging", str(evolution['summary']['converging_toward_mrbeast']), "of panel channels"),
                                ("Avg Slope", f"+{evolution['summary']['avg_slope']:.3f}/yr", None)])

    # Slides 20-23: Case Studies
    case_studies = [
        ("Danny Duncan", "Fastest converger — surpassed MrBeast's average by 2022",
         [("Start (2015)", "2.00", None), ("End (2025)", "6.67", "Exceeds MrBeast avg"), ("Slope", "+0.551/yr", "Fastest"), ("Peak", "6.67", "2025")]),
        ("ZHC", "Peaked at 7.40 in 2024 — highest single-year score of any panel channel",
         [("Start (2017)", "1.47", None), ("End (2025)", "4.53", None), ("Slope", "+0.434/yr", "2nd fastest"), ("Peak", "7.40", "2024")]),
        ("FaZe Rug", "Slow start, rapid adoption — jumped from 3.1 to 6.6 in three years",
         [("Start (2015)", "3.33", None), ("End (2025)", "6.47", "Exceeds MrBeast avg"), ("Slope", "+0.355/yr", None), ("2021→2024", "3.1→6.6", "Rapid jump")]),
        ("Sidemen", "Steady climb from 1.3 to 5.1 — now approaching MrBeast's average",
         [("Start (2016)", "1.29", None), ("End (2025)", "5.07", "94% of MrBeast avg"), ("Slope", "+0.314/yr", None), ("Gap", "0.35 pts", "from MrBeast")]),
    ]
    for ch_name, subtitle, stat_cards in case_studies:
        ch_data = evolution.get('channels', {}).get(ch_name, {})
        years = sorted(ch_data.get('years', {}).keys())
        scores = [ch_data['years'][y]['mean_score'] for y in years]
        slide_case_study(prs, ch_name, subtitle, years, scores, stat_cards)

    # Slide 24: Weighted Feature Importance
    weights = sorted(weighted_weights['weights'].items(), key=lambda x: -x[1])
    slide_chart_bar(prs, "Weighted Feature Importance",
                    "Data-derived weights: how discriminative is each feature?",
                    [w[0].replace('_', ' ') for w in weights],
                    [w[1] for w in weights],
                    series_name="Weight")

    # Slide 25: Statistical Validation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Statistical Validation")
    t = convergence.get('ttest', {})
    a = convergence.get('anova', {})
    r = convergence.get('linear_regression', {})
    d = convergence.get('cohens_d', {})
    cards = [
        ("Welch's t-test", f"t = {t.get('t_statistic', 0):.2f}", f"p = {t.get('p_value', 0):.2e}"),
        ("ANOVA", f"F = {a.get('f_statistic', 0):.2f}", f"p = {a.get('p_value', 0):.2e}"),
        ("Regression", f"slope = +{r.get('slope', 0):.3f}/yr", f"p = {r.get('p_value', 0):.2e}"),
        ("Cohen's d", f"{d.get('d', 0):.2f}", d.get('interpretation', '')),
    ]
    for i, (label, val, sub) in enumerate(cards):
        add_stat_card(slide, label, val, Inches(0.3 + i * 2.4), Inches(2), sub=sub)

    # CI line chart
    ci = convergence.get('year_confidence_intervals', {})
    ci_years = [g for g in YEAR_GROUPS if g in ci]
    ci_means = [ci[g]['mean'] for g in ci_years]
    if ci_years:
        add_line_chart(slide, ci_years, {"Mean Likeness": ci_means},
                       left=Inches(0.5), top=Inches(3.5), width=Inches(9), height=Inches(3.5))

    # Slide 26: Title Convergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Title Convergence")
    add_bullets(slide, [
        "Title likeness +14.6% (vs +27% for thumbnails)",
        "First-person framing reached parity (25.4% vs 25.9%)",
        "Money references: 5x increase (2.3% → 11.8%)",
        "Challenge framing largely unadopted (13.3% vs 39.8%)",
        "Numeric hooks still far behind (32.6% vs 67.3%)",
    ], left=Inches(0.5), top=Inches(1.8), width=Inches(4.5), font_size=13)

    t_values = [title_data['groups'].get(g, {}).get('mean_score', 0) for g in YEAR_GROUPS]
    add_bar_chart(slide, YEAR_GROUPS, t_values, title="Mean Title Likeness by Year",
                  left=Inches(5), top=Inches(1.8), width=Inches(4.8), height=Inches(3.5),
                  series_name="Title Likeness")

    # Slide 27: Diffusion of Innovation
    slide_27_diffusion(prs)

    # Slide 28: Limitations & Conclusions
    slide_28_conclusions(prs)

    # Save
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"✓ Saved to {OUTPUT}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
